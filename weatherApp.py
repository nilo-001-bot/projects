from pptx import Presentation

# Create a new PowerPoint presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Interaction Models in Software Engineering"
subtitle.text = "An Overview of System Models and Interactions"

# Slide 2: Introduction to Interaction Models
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "• Interaction models describe how systems interact internally and externally.\n"
    "• They capture:\n"
    "  - User interactions (inputs/outputs)\n"
    "  - System-to-system interactions\n"
    "  - Component interactions\n"
    "• Benefits include identifying requirements, resolving communication issues, and ensuring system dependability."
)

# Slide 3: Types of Interaction Models
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Types of Interaction Models"
content.text = (
    "1. User Interaction:\n"
    "   - Focuses on how users interact with the system.\n"
    "   - Examples: input forms, UI workflows.\n\n"
    "2. System-to-System Interaction:\n"
    "   - Explores data exchanges between different systems.\n"
    "   - Examples: API calls, service integrations.\n\n"
    "3. Component Interaction:\n"
    "   - Examines collaboration among internal components.\n"
    "   - Examples: microservices or modular interactions."
)

# Slide 4: Use Case Modeling
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Use Case Modeling"
content.text = (
    "• Represents tasks involving external interactions.\n"
    "• Visualized using ellipses (use cases) and stick figures (actors).\n"
    "• Examples: System login, data upload processes.\n"
    "• Benefits:\n"
    "  - Captures user expectations.\n"
    "  - Simplifies requirements documentation."
)

# Slide 5: Sequence Diagrams
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Sequence Diagrams"
content.text = (
    "• Models interactions among components, users, and systems.\n"
    "• Key features:\n"
    "  - Objects and actors at the top.\n"
    "  - Interactions as arrows.\n"
    "  - Lifelines as vertical dotted lines.\n"
    "• Highlights sequence and timing of interactions."
)

# Slide 6: Benefits of Interaction Models
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Benefits of Interaction Models"
content.text = (
    "• Clarify requirements and expectations.\n"
    "• Improve communication between stakeholders.\n"
    "• Ensure system performance and dependability.\n"
    "• Simplify testing and validation processes."
)

# Slide 7: Conclusion
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "• Interaction models are vital for system design and analysis.\n"
    "• Use case models provide high-level overviews.\n"
    "• Sequence diagrams offer detailed views of component interactions.\n"
    "• Together, they ensure comprehensive system documentation."
)

# Save the presentation
file_path = "C:/Users/Hp/Downloads/Interaction_Models_Presentation_Final.pptx"
presentation.save(file_path)

file_path

